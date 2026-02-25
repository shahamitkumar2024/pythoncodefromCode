"""Tool definitions for all agents in the Financial Intelligence system."""

import os
import re
from typing import Optional

import requests
from bs4 import BeautifulSoup
from langchain_core.tools import tool


# ---------------------------------------------------------------------------
# Web Downloader Tools
# ---------------------------------------------------------------------------

@tool
def search_company_report(company_name: str) -> str:
    """Search for the last two years of annual reports and P&L statements for a company.

    Uses Tavily if TAVILY_API_KEY is set; falls back to DuckDuckGo HTML scrape.

    Args:
        company_name: Name of the target company (e.g. "Microsoft").

    Returns:
        Top search results (title + URL + snippet) covering FY2024 and FY2023.
    """
    query = f"{company_name} annual report P&L investor relations 2024 2023"
    tavily_key = os.getenv("TAVILY_API_KEY")

    if tavily_key:
        try:
            from tavily import TavilyClient  # type: ignore[import-untyped]

            client = TavilyClient(api_key=tavily_key)
            results = client.search(query=query, max_results=5)
            formatted = []
            for r in results.get("results", []):
                formatted.append(
                    f"Title: {r.get('title', '')}\n"
                    f"URL: {r.get('url', '')}\n"
                    f"Snippet: {r.get('content', '')[:300]}\n"
                )
            return "\n---\n".join(formatted) if formatted else "No results found."
        except Exception as e:
            return f"Tavily search failed: {e}. Try fetch_page_content with a direct URL."

    # Fallback: DuckDuckGo HTML search (no API key needed)
    headers = {"User-Agent": "Mozilla/5.0 (compatible; FinancialIntelAgent/1.0)"}
    search_url = f"https://html.duckduckgo.com/html/?q={requests.utils.quote(query)}"
    try:
        resp = requests.get(search_url, headers=headers, timeout=15)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")
        results = []
        for result in soup.select(".result")[:5]:
            title_el = result.select_one(".result__title")
            url_el = result.select_one(".result__url")
            snippet_el = result.select_one(".result__snippet")
            title = title_el.get_text(strip=True) if title_el else ""
            url = url_el.get_text(strip=True) if url_el else ""
            snippet = snippet_el.get_text(strip=True) if snippet_el else ""
            if title:
                results.append(f"Title: {title}\nURL: {url}\nSnippet: {snippet[:300]}")
        return "\n---\n".join(results) if results else "No results found via fallback search."
    except Exception as e:
        return (
            f"Fallback search failed: {e}. "
            "Please provide a direct investor relations URL to fetch_page_content."
        )


@tool
def fetch_page_content(url: str) -> str:
    """Fetch and extract readable text content from a web page or PDF.

    Automatically detects PDF responses and extracts text using pypdf.
    Falls back to BeautifulSoup HTML parsing for regular web pages.

    Args:
        url: URL of an investor relations page, annual report HTML, or PDF.

    Returns:
        Cleaned plain-text content (up to ~50,000 chars).
    """
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/120.0.0.0 Safari/537.36"
        )
    }
    try:
        resp = requests.get(url, headers=headers, timeout=30)
        resp.raise_for_status()
    except requests.RequestException as e:
        return f"Failed to fetch {url}: {e}"

    # ── PDF detection ──────────────────────────────────────────────────────
    content_type = resp.headers.get("content-type", "").lower()
    is_pdf = "pdf" in content_type or url.lower().split("?")[0].endswith(".pdf")

    if is_pdf:
        try:
            import io
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(resp.content))
            pages_text: list[str] = []
            for page in reader.pages[:60]:          # cap at 60 pages
                text = page.extract_text() or ""
                if text.strip():
                    pages_text.append(text)

            extracted = "\n\n".join(pages_text)
            if extracted.strip():
                return _clean_text(extracted)[:50_000]

            return (
                "PDF downloaded but no text could be extracted "
                "(likely a scanned/image-only PDF). "
                "Please provide the HTML investor relations page URL instead."
            )
        except ImportError:
            return (
                "pypdf is not installed — cannot read PDF files. "
                "Run: pip install pypdf\n"
                "Alternatively, provide the HTML investor relations page URL."
            )
        except Exception as e:
            return (
                f"PDF extraction failed: {e}. "
                "Try the HTML investor relations page URL instead."
            )

    # ── HTML parsing ───────────────────────────────────────────────────────
    soup = BeautifulSoup(resp.text, "lxml")

    for tag in soup(["script", "style", "nav", "footer", "header", "aside", "form"]):
        tag.decompose()

    for selector in ["main", "article", "#content", ".content", "#main-content"]:
        main = soup.select_one(selector)
        if main:
            text = main.get_text(separator="\n", strip=True)
            return _clean_text(text)[:50_000]

    text = soup.get_text(separator="\n", strip=True)
    return _clean_text(text)[:50_000]


@tool
def extract_financial_tables(page_content: str) -> str:
    """Extract key financial metrics from page text for both FY2024 and FY2023.

    Args:
        page_content: Raw text content from a financial report page.

    Returns:
        Structured summary of financial metrics found in the text.
    """
    patterns = {
        "Revenue": r"(?:total\s+)?revenue[s]?\s*[:\-]?\s*\$?([\d,\.]+\s*(?:billion|million|B|M)?)",
        "Net Income": r"net\s+(?:income|profit|earnings)\s*[:\-]?\s*\$?([\d,\.]+\s*(?:billion|million|B|M)?)",
        "EBITDA": r"EBITDA\s*[:\-]?\s*\$?([\d,\.]+\s*(?:billion|million|B|M)?)",
        "Operating Income": r"operating\s+(?:income|profit)\s*[:\-]?\s*\$?([\d,\.]+\s*(?:billion|million|B|M)?)",
        "EPS": r"(?:diluted\s+)?(?:EPS|earnings\s+per\s+share)\s*[:\-]?\s*\$?([\d,\.]+)",
        "Gross Profit": r"gross\s+profit\s*[:\-]?\s*\$?([\d,\.]+\s*(?:billion|million|B|M)?)",
        "Free Cash Flow": r"free\s+cash\s+flow\s*[:\-]?\s*\$?([\d,\.]+\s*(?:billion|million|B|M)?)",
        "R&D Spend": r"research\s+(?:and|&)\s+development\s*[:\-]?\s*\$?([\d,\.]+\s*(?:billion|million|B|M)?)",
    }

    found: dict[str, list[str]] = {}
    content_lower = page_content.lower()
    for metric, pattern in patterns.items():
        matches = re.findall(pattern, content_lower, re.IGNORECASE)
        if matches:
            found[metric] = list(dict.fromkeys(matches[:4]))  # keep up to 4 (may be 2 years × 2)

    if not found:
        return (
            "No structured financial tables could be automatically extracted. "
            "The analyst should parse the raw report text directly. "
            "Look for sections titled 'Consolidated Income Statement', "
            "'Financial Highlights', or 'Key Performance Indicators'."
        )

    lines = ["## Extracted Financial Metrics (may include multiple years)\n"]
    for metric, values in found.items():
        lines.append(f"- **{metric}**: {' | '.join(values)}")

    lines.append(
        "\n_Note: Values extracted via pattern matching — first value is likely most recent year. "
        "Analyst should verify and pair each value to the correct fiscal year._"
    )
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Financial Analyst Tools
# ---------------------------------------------------------------------------

@tool
def analyze_financial_health(raw_report: str) -> str:
    """Analyse financial health from two years of report text.

    Returns a checklist of what multi-year data is available in the raw report.

    Args:
        raw_report: The full raw text of the financial report(s).

    Returns:
        Checklist of available data points with two-year coverage indicators.
    """
    word_count = len(raw_report.split())
    has_revenue = bool(re.search(r"revenue", raw_report, re.IGNORECASE))
    has_profit = bool(re.search(r"net\s+(income|profit)", raw_report, re.IGNORECASE))
    has_ebitda = bool(re.search(r"ebitda", raw_report, re.IGNORECASE))
    has_guidance = bool(re.search(r"guidance|outlook|forecast", raw_report, re.IGNORECASE))
    has_debt = bool(re.search(r"debt|borrowing|liabilit", raw_report, re.IGNORECASE))
    has_cashflow = bool(re.search(r"cash\s+flow|free\s+cash", raw_report, re.IGNORECASE))
    has_2024 = bool(re.search(r"2024|fy24|fy\s*2024", raw_report, re.IGNORECASE))
    has_2023 = bool(re.search(r"2023|fy23|fy\s*2023", raw_report, re.IGNORECASE))

    checklist = [
        f"- Report length: ~{word_count} words",
        f"- FY2024 data present: {'YES' if has_2024 else 'NO — check for alternative year labels'}",
        f"- FY2023 data present: {'YES ✓ (two-year comparison available)' if has_2023 else 'NO — fetch prior year report'}",
        f"- Revenue data: {'YES' if has_revenue else 'NO'}",
        f"- Profit data: {'YES' if has_profit else 'NO'}",
        f"- EBITDA data: {'YES' if has_ebitda else 'NO'}",
        f"- Cash flow data: {'YES' if has_cashflow else 'NO'}",
        f"- Debt/liability data: {'YES' if has_debt else 'NO'}",
        f"- Forward guidance: {'YES' if has_guidance else 'NO'}",
    ]

    return (
        "## Two-Year Financial Health Analysis Checklist\n"
        + "\n".join(checklist)
        + "\n\n"
        "Extract actual values for BOTH years and compute:\n"
        "1. Revenue: FY2024 value, FY2023 value, YoY growth %\n"
        "2. Net Profit Margin %: both years\n"
        "3. EBITDA Margin %: both years (if available)\n"
        "4. Free Cash Flow trend: FY2023 → FY2024 (improving/stable/declining)\n"
        "5. R&D / capex investment: both years\n"
        "6. Overall health trend: IMPROVING / STABLE / DECLINING"
    )


@tool
def extract_strategic_initiatives(raw_report: str) -> str:
    """Extract strategic priorities from financial report text.

    Detects transformation themes and flags whether they appear in both years.

    Args:
        raw_report: The full raw text of the financial report(s).

    Returns:
        Detected strategic themes with year coverage.
    """
    themes = {
        "Digital Transformation": ["digital transformation", "digitization", "digitalization", "digital strategy"],
        "Cloud Migration": ["cloud", "aws", "azure", "google cloud", "hybrid cloud", "cloud-first"],
        "AI & Automation": ["artificial intelligence", "machine learning", "automation", "ai", "genai", "generative ai", "robotic process"],
        "Cost Optimization": ["cost reduction", "efficiency", "optimization", "restructuring", "cost savings"],
        "Cybersecurity": ["cybersecurity", "security", "data protection", "zero trust", "compliance"],
        "Data & Analytics": ["data analytics", "big data", "data platform", "analytics", "data-driven"],
        "Customer Experience": ["customer experience", "cx", "omnichannel", "customer journey", "personalization"],
        "Supply Chain": ["supply chain", "logistics", "procurement", "inventory"],
        "ESG / Sustainability": ["sustainability", "esg", "carbon", "net zero", "green"],
        "M&A / Expansion": ["acquisition", "merger", "expansion", "market entry", "growth strategy"],
    }

    report_lower = raw_report.lower()
    found_themes: list[str] = []
    for theme, keywords in themes.items():
        hits = [kw for kw in keywords if kw in report_lower]
        if hits:
            found_themes.append(f"- **{theme}** (keywords: {', '.join(hits[:3])})")

    if not found_themes:
        return (
            "No explicit strategic theme keywords detected. "
            "Review the MD&A and CEO letter sections manually."
        )

    return (
        "## Detected Strategic Themes\n"
        + "\n".join(found_themes)
        + "\n\n"
        "For each theme, extract 3-5 specific initiatives noting:\n"
        "- Whether the initiative appears in both FY2024 and FY2023 reports (ongoing) or is new\n"
        "- Budget/investment mentioned and any timeline\n"
        "- Progress/changes between the two years (e.g. 'Phase 1 complete in FY2023, Phase 2 in FY2024')"
    )


@tool
def generate_opportunity_map(analysis: str) -> str:
    """Map identified company initiatives to TCS service offering categories.

    Args:
        analysis: The structured financial and strategic analysis text.

    Returns:
        A Markdown opportunity map linking company needs to TCS capabilities.
    """
    tcs_services = {
        "Digital Transformation": [
            "TCS Digital Transformation Services",
            "TCS BaNCS (BFSI)",
            "TCS iON (SME platforms)",
            "Consulting & Advisory",
        ],
        "Cloud Migration": [
            "TCS Cloud Services (AWS/Azure/GCP)",
            "TCS HOBS (cloud-native platforms)",
            "Infrastructure Managed Services",
        ],
        "AI & Automation": [
            "TCS AI.Cloud",
            "Ignio AIOps (Digitate)",
            "TCS Robotic Process Automation",
            "Generative AI COE Services",
        ],
        "Cost Optimization": [
            "Application Outsourcing (AO)",
            "Business Process Services (BPS)",
            "TCS MasterCraft (automated testing & dev)",
        ],
        "Cybersecurity": [
            "TCS Cyber Security Practice",
            "Managed Security Operations Center (MSOC)",
            "Zero Trust Architecture Consulting",
        ],
        "Data & Analytics": [
            "TCS Data & Analytics Services",
            "TCS HOBS Data Platform",
            "Enterprise Data Warehousing / Lakehouse",
        ],
        "Customer Experience": [
            "TCS Customer Intelligence & Insights (CI&I)",
            "TCS OmniStore (Retail)",
            "Salesforce & CRM Implementation",
        ],
        "Supply Chain": [
            "TCS Supply Chain Operations",
            "SAP SCM Implementation",
            "Digital Twin & IoT Solutions",
        ],
        "ESG / Sustainability": [
            "TCS ESG Framework & Reporting",
            "Green IT Consulting",
            "Carbon Footprint Management Platforms",
        ],
    }

    lines = ["## TCS Opportunity Map\n"]
    for theme, services in tcs_services.items():
        if theme.lower() in analysis.lower() or any(
            s.lower() in analysis.lower() for s in theme.lower().split()
        ):
            lines.append(f"### {theme}")
            for svc in services:
                lines.append(f"  - {svc}")
            lines.append("")

    if len(lines) == 1:
        lines.append("Review the analysis and map initiatives to TCS service lines manually.")

    lines.append(
        "_Use this map for the TCS Service Alignment slides. "
        "Customise to the specific company context._"
    )
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Presentation Tools
# ---------------------------------------------------------------------------

@tool
def format_tcs_value_props(initiatives: str) -> str:
    """Format TCS value propositions for specific company initiatives.

    Args:
        initiatives: Comma-separated list of company strategic initiative names.

    Returns:
        Formatted TCS value proposition bullets for each initiative.
    """
    props = {
        "cloud": (
            "TCS cloud-agnostic delivery (AWS/Azure/GCP certified), "
            "500+ cloud migrations delivered, average 30% TCO reduction"
        ),
        "ai": (
            "TCS AI.Cloud accelerators, Ignio AIOps reducing MTTR by 40%, "
            "Generative AI COE with 10,000+ trained practitioners"
        ),
        "digital": (
            "TCS Digital Platform Studios, industry-specific accelerators, "
            "100+ digital transformation programs in Fortune 500"
        ),
        "cost": (
            "TCS BPS delivering 25-40% cost reduction, "
            "MasterCraft automation reducing testing effort by 60%"
        ),
        "security": (
            "TCS Cyber Security practice, MSOC covering 50+ countries, "
            "SOC 2 / ISO 27001 certified delivery centers"
        ),
        "data": (
            "TCS Data & Analytics practice, proprietary data lakehouse accelerators, "
            "GDPR/CCPA compliant data management"
        ),
    }

    result_lines = ["## TCS Value Propositions by Initiative\n"]
    for keyword, prop in props.items():
        if keyword in initiatives.lower():
            result_lines.append(f"**{keyword.title()} Initiative:**")
            result_lines.append(f"  {prop}\n")

    if len(result_lines) == 1:
        result_lines.append(
            "Map initiatives to TCS capabilities using the opportunity map. "
            "Highlight TCS global scale, IP-led solutions, and industry expertise."
        )

    return "\n".join(result_lines)


# ---------------------------------------------------------------------------
# PowerPoint builder helpers
# ---------------------------------------------------------------------------

def _clean_text(text: str) -> str:
    """Remove excessive whitespace from extracted text."""
    lines = [line.strip() for line in text.splitlines()]
    non_empty = [line for line in lines if line]
    return "\n".join(non_empty)


def _clean_markdown(text: str) -> str:
    """Strip markdown syntax for plain-text display in PowerPoint."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)
    return text.strip()


def _parse_slides(markdown: str) -> list[dict]:
    """Split markdown into slide dicts: title, text_lines, table_rows, has_table."""
    sections = re.split(r"\n\s*---+\s*\n", markdown)
    slides = []

    for section in sections:
        section = section.strip()
        if not section:
            continue

        lines = section.splitlines()
        title = ""
        body_start = 0

        for i, line in enumerate(lines):
            clean = line.strip()
            if re.match(r"^#{1,3}\s", clean):
                title = re.sub(r"^#+\s*", "", clean)
                title = re.sub(r"^Slide\s+\d+[:\.\)]\s*", "", title, flags=re.IGNORECASE)
                body_start = i + 1
                break

        if not title:
            for i, line in enumerate(lines):
                if line.strip():
                    title = line.strip().lstrip("*_#").strip()
                    body_start = i + 1
                    break

        body_lines = lines[body_start:]
        table_lines = [l for l in body_lines if re.match(r"^\s*\|", l)]
        text_lines = [l for l in body_lines if not re.match(r"^\s*\|", l)]
        table_rows = _parse_markdown_table(table_lines) if len(table_lines) >= 2 else []

        slides.append({
            "title": title,
            "text_lines": text_lines,
            "table_rows": table_rows,
            "has_table": bool(table_rows),
        })

    return slides


def _parse_markdown_table(lines: list[str]) -> list[list[str]]:
    """Parse markdown table lines into row-lists, skipping separator rows."""
    rows = []
    for line in lines:
        line = line.strip()
        if not line.startswith("|"):
            continue
        if re.match(r"^\|[\s\-\|:]+\|$", line):
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        if any(c for c in cells):
            rows.append(cells)
    return rows


def _apply_tcs_theme(prs: object) -> None:
    from pptx.util import Inches
    prs.slide_width = Inches(13.33)   # type: ignore[attr-defined]
    prs.slide_height = Inches(7.5)    # type: ignore[attr-defined]


def _add_cover_slide(prs: object, company_name: str) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    TCS_DARK = RGBColor(0, 40, 80)
    TCS_BLUE = RGBColor(0, 84, 166)
    TCS_CYAN = RGBColor(0, 169, 224)
    WHITE    = RGBColor(255, 255, 255)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # type: ignore[attr-defined]

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = TCS_DARK

    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.18)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = TCS_CYAN
    stripe.line.fill.background()

    sub = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.33), Inches(0.8))
    sub.text_frame.word_wrap = True
    p = sub.text_frame.paragraphs[0]
    p.text = "TCS Business Intelligence"
    p.font.size = Pt(24)
    p.font.color.rgb = TCS_CYAN
    p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.8))
    title_box.text_frame.word_wrap = True
    p2 = title_box.text_frame.paragraphs[0]
    p2.text = company_name
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    sub2 = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.33), Inches(0.6))
    p3 = sub2.text_frame.paragraphs[0]
    p3.text = "Business Opportunity Analysis  |  Prepared by TCS Business Development"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TCS_CYAN
    p3.alignment = PP_ALIGN.CENTER

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.33), Inches(0.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TCS_BLUE
    bar.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(12.73), Inches(0.35))
    pf = footer.text_frame.paragraphs[0]
    pf.text = "CONFIDENTIAL — For TCS Internal Use Only"
    pf.font.size = Pt(9)
    pf.font.color.rgb = WHITE
    pf.alignment = PP_ALIGN.CENTER


def _add_title_bar(slide: object, title: str) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    TCS_BLUE = RGBColor(0, 84, 166)
    WHITE    = RGBColor(255, 255, 255)

    bar = slide.shapes.add_shape(  # type: ignore[attr-defined]
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TCS_BLUE
    bar.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.73), Inches(0.9))  # type: ignore[attr-defined]
    txb.text_frame.word_wrap = True
    p = txb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE


def _add_footer_bar(slide: object) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    TCS_CYAN  = RGBColor(0, 169, 224)
    DARK_GRAY = RGBColor(90, 90, 90)

    bar = slide.shapes.add_shape(  # type: ignore[attr-defined]
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.33), Inches(0.32)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TCS_CYAN
    bar.line.fill.background()

    ft = slide.shapes.add_textbox(Inches(0.3), Inches(7.2), Inches(12.73), Inches(0.28))  # type: ignore[attr-defined]
    p = ft.text_frame.paragraphs[0]
    p.text = "CONFIDENTIAL — TCS Business Development"
    p.font.size = Pt(8)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT


def _add_content_slide(prs: object, slide_data: dict) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE

    TCS_BLUE  = RGBColor(0, 84, 166)
    TCS_DARK  = RGBColor(0, 40, 80)
    DARK_TEXT = RGBColor(35, 35, 35)
    MID_GRAY  = RGBColor(80, 80, 80)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # type: ignore[attr-defined]
    _add_title_bar(slide, slide_data["title"])

    txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.7))
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    first = True
    for raw_line in slide_data["text_lines"]:
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            continue
        if re.match(r"^\*Prepared by", stripped):
            continue
        if re.match(r"^#{1,3}\s", stripped):
            continue

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        num_match     = re.match(r"^(\d+)[\.)\s]\s+(.+)", stripped)
        bullet_match  = re.match(r"^[\-\*\•]\s+(.+)", stripped)
        check_match   = re.match(r"^[\-\*]\s+\[[ xX]\]\s+(.+)", stripped)
        head3_match   = re.match(r"^###\s+(.+)", stripped)
        bold_head_match  = re.match(r"^\*\*(.+?)\*\*\s*$", stripped)
        bold_label_match = re.match(r"^\*\*(.+?)\*\*[:\s]+(.+)", stripped)

        if check_match:
            p.text = f"☐  {_clean_markdown(check_match.group(1))}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(3)
        elif bullet_match:
            p.text = f"  •  {_clean_markdown(bullet_match.group(1))}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(2)
        elif num_match:
            p.text = f"  {num_match.group(1)}.  {_clean_markdown(num_match.group(2))}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(3)
        elif head3_match:
            p.text = _clean_markdown(head3_match.group(1))
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = TCS_BLUE
            p.space_before = Pt(8)
        elif bold_head_match:
            p.text = _clean_markdown(bold_head_match.group(1))
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = TCS_DARK
            p.space_before = Pt(6)
        elif bold_label_match:
            label = _clean_markdown(bold_label_match.group(1))
            value = _clean_markdown(bold_label_match.group(2))
            run_label = p.add_run()
            run_label.text = f"{label}:  "
            run_label.font.bold = True
            run_label.font.size = Pt(13)
            run_label.font.color.rgb = TCS_DARK
            run_value = p.add_run()
            run_value.text = value
            run_value.font.size = Pt(13)
            run_value.font.color.rgb = DARK_TEXT
            p.space_before = Pt(3)
        else:
            p.text = _clean_markdown(stripped)
            p.font.size = Pt(13)
            p.font.color.rgb = MID_GRAY
            p.space_before = Pt(2)

    _add_footer_bar(slide)


def _add_table_slide(prs: object, slide_data: dict) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    TCS_BLUE    = RGBColor(0, 84, 166)
    TCS_ALT_ROW = RGBColor(235, 243, 255)
    WHITE       = RGBColor(255, 255, 255)
    DARK_TEXT   = RGBColor(35, 35, 35)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # type: ignore[attr-defined]
    _add_title_bar(slide, slide_data["title"])

    table_rows = slide_data["table_rows"]
    if not table_rows:
        _add_footer_bar(slide)
        return

    y_offset = 1.28
    intro_lines = [
        l.strip() for l in slide_data["text_lines"]
        if l.strip() and not re.match(r"^#{1,3}\s", l.strip())
    ][:2]

    if intro_lines:
        txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.25), Inches(12.33), Inches(0.65))
        txb.text_frame.word_wrap = True
        for i, intro in enumerate(intro_lines):
            p = txb.text_frame.paragraphs[0] if i == 0 else txb.text_frame.add_paragraph()
            p.text = _clean_markdown(intro)
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
        y_offset = 1.95

    n_rows = len(table_rows)
    n_cols = max(len(r) for r in table_rows)
    table_rows = [r + [""] * (n_cols - len(r)) for r in table_rows]

    row_height = Inches(0.52)
    table_height = min(row_height * n_rows, Inches(7.18) - Inches(y_offset) - Inches(0.4))

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.4), Inches(y_offset),
        Inches(12.5), table_height,
    ).table

    for r_idx, row_data in enumerate(table_rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = _clean_markdown(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TCS_BLUE
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TCS_ALT_ROW
                p.font.color.rgb = DARK_TEXT
            else:
                p.font.color.rgb = DARK_TEXT

    _add_footer_bar(slide)


@tool
def save_as_powerpoint(content: str, company_name: str) -> str:
    """Convert a Markdown presentation into a professional PowerPoint .pptx file.

    Parses slide sections separated by '---', renders bullet content and
    markdown tables as native PPTX elements, and applies TCS brand colours.

    Args:
        content: Full Markdown presentation text (slides separated by ---).
        company_name: Used for the cover slide title and output filename.

    Returns:
        Absolute path to the saved .pptx file.
    """
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"

    from pptx import Presentation

    prs = Presentation()
    _apply_tcs_theme(prs)
    _add_cover_slide(prs, company_name)

    for slide_data in _parse_slides(content):
        if slide_data["has_table"]:
            _add_table_slide(prs, slide_data)
        else:
            _add_content_slide(prs, slide_data)

    safe_name = re.sub(r"[^\w\-]", "_", company_name.lower())
    filename = f"presentation_{safe_name}.pptx"
    output_path = os.path.join(os.getcwd(), filename)
    prs.save(output_path)

    return f"PowerPoint saved: {output_path}  ({len(prs.slides)} slides)"


# ---------------------------------------------------------------------------
# Reviewer Tools
# ---------------------------------------------------------------------------

@tool
def verify_financial_facts(presentation: str, raw_report: str) -> str:
    """Cross-check financial figures in the presentation against the raw report.

    Args:
        presentation: The Markdown presentation content.
        raw_report: The original raw financial report text.

    Returns:
        A verification report listing each figure checked and its status.
    """
    pres_figures = re.findall(
        r"\$[\d,\.]+\s*(?:billion|million|B|M|b|m)?|£[\d,\.]+\s*(?:billion|million|B|M|b|m)?|\d+\.?\d*\s*%",
        presentation,
        re.IGNORECASE,
    )

    report_lower = raw_report.lower()
    lines = ["## Financial Fact Verification\n"]
    verified = []
    unverified = []

    for figure in dict.fromkeys(pres_figures):
        normalized = re.sub(r"[\$£,]", "", figure).strip().lower()
        core_number = re.sub(r"[^\d\.]", "", normalized)
        if core_number and core_number in report_lower:
            verified.append(figure)
        else:
            unverified.append(figure)

    if verified:
        lines.append("### Verified Figures")
        for fig in verified:
            lines.append(f"  - {fig} ✓")
        lines.append("")

    if unverified:
        lines.append("### Figures Requiring Manual Verification")
        for fig in unverified:
            lines.append(f"  - [UNVERIFIED: {fig}] — not found verbatim in raw report")
        lines.append("")

    if not pres_figures:
        lines.append(
            "No specific financial figures found in presentation. "
            "Ensure the presentation contains actual data for both FY2024 and FY2023."
        )

    return "\n".join(lines)


@tool
def check_tcs_claims(presentation: str) -> str:
    """Validate TCS capability claims in the presentation are realistic.

    Args:
        presentation: The Markdown presentation content.

    Returns:
        A list of flagged overclaims or vague statements.
    """
    vague_patterns = [
        (r"\bbest[\- ]in[\- ]class\b", "vague superlative — add specific metric"),
        (r"\bworld[\- ]class\b", "vague superlative — add specific proof point"),
        (r"\bunmatched\b", "unsubstantiated claim — add comparative data"),
        (r"\bseamlessly\b", "marketing fluff — replace with specific capability"),
        (r"\bsynergies\b", "vague — specify what synergies"),
        (r"\b100%\s+guarantee\b", "overclaim — use SLA-based language"),
        (r"\bguarantee[sd]?\b", "risky commitment language — use 'target' or 'SLA'"),
        (r"\binstantly\b", "unrealistic — add realistic timeline"),
    ]

    issues = []
    for pattern, advice in vague_patterns:
        match = re.search(pattern, presentation, re.IGNORECASE)
        if match:
            issues.append(f"  - [OVERCLAIM: '{match.group()}'] — {advice}")

    tcs_products = ["tcs", "ignio", "bancs", "mastercraft", "ion", "ai.cloud"]
    has_tcs_ref = any(p in presentation.lower() for p in tcs_products)

    lines = ["## TCS Claims Review\n"]
    if issues:
        lines.append("### Flagged Language")
        lines.extend(issues)
        lines.append("")
    else:
        lines.append("### Language Review: PASS — No overclaims detected\n")

    if not has_tcs_ref:
        lines.append(
            "### [MISSING: TCS-specific product references] — "
            "Add named TCS products/platforms (e.g., TCS BaNCS, Ignio, AI.Cloud)\n"
        )
    else:
        lines.append("### TCS Product References: PRESENT ✓\n")

    return "\n".join(lines)


@tool
def generate_review_report(
    fact_verification: str,
    claims_check: str,
    presentation: str,
) -> str:
    """Compile all review findings into a structured review report with PASS/NEEDS_REVISION verdict.

    Args:
        fact_verification: Output from verify_financial_facts.
        claims_check: Output from check_tcs_claims.
        presentation: The presentation being reviewed.

    Returns:
        Final structured review report.
    """
    required_sections = [
        "Executive Summary",
        "Financial Health",
        "Strategic Priorities",
        "TCS Service Alignment",
        "Engagement Model",
        "Next Steps",
    ]

    missing_sections = [
        s for s in required_sections if s.lower() not in presentation.lower()
    ]

    has_two_years = bool(
        re.search(r"fy2024|fy\s*2024|2024", presentation, re.IGNORECASE)
        and re.search(r"fy2023|fy\s*2023|2023", presentation, re.IGNORECASE)
    )

    has_unverified = "[UNVERIFIED" in fact_verification
    has_overclaims = "[OVERCLAIM" in claims_check
    has_missing_tcs = "[MISSING: TCS" in claims_check
    issues_found = has_unverified or has_overclaims or missing_sections or has_missing_tcs or not has_two_years

    verdict = "NEEDS_REVISION" if issues_found else "PASS"

    report = [
        "# Presentation Review Report",
        "",
        "## 1. Financial Fact Verification",
        fact_verification,
        "",
        "## 2. TCS Claims Review",
        claims_check,
        "",
        "## 3. Two-Year Coverage Check",
        f"FY2024 data present: {'YES' if re.search(r'2024', presentation) else 'NO'}",
        f"FY2023 data present: {'YES' if re.search(r'2023', presentation) else 'NO'}",
        f"Two-year comparison: {'COMPLETE ✓' if has_two_years else '[MISSING_YEAR: presentation must include FY2023 and FY2024 side-by-side]'}",
        "",
        "## 4. Structure Check",
    ]

    if missing_sections:
        report.append(f"Missing required sections: {', '.join(missing_sections)}")
        for sec in missing_sections:
            report.append(f"  - [MISSING_SECTION: {sec}]")
    else:
        report.append("All required sections present ✓")

    report += [
        "",
        "## Summary",
        f"- Two-year data: {'COMPLETE' if has_two_years else 'MISSING FY2023'}",
        f"- Unverified figures: {'YES' if has_unverified else 'None'}",
        f"- Overclaims: {'YES' if has_overclaims else 'None'}",
        f"- Missing sections: {', '.join(missing_sections) if missing_sections else 'None'}",
        "",
        f"**VERDICT: {verdict}**",
    ]

    if verdict == "NEEDS_REVISION":
        report.append("\n### Required Corrections:")
        if not has_two_years:
            report.append("- Add FY2023 data alongside FY2024 in the Financial Health slide with YoY comparisons")
        if has_unverified:
            report.append("- Replace [UNVERIFIED] figures with verified data from the raw report")
        if has_overclaims:
            report.append("- Replace vague/hyperbolic language with specific, data-backed claims")
        if missing_sections:
            report.append(f"- Add missing sections: {', '.join(missing_sections)}")
        if has_missing_tcs:
            report.append("- Add specific TCS product/platform references")

    return "\n".join(report)


# ---------------------------------------------------------------------------
# Tool groups
# ---------------------------------------------------------------------------

WEB_DOWNLOADER_TOOLS = [search_company_report, fetch_page_content, extract_financial_tables]
ANALYST_TOOLS = [analyze_financial_health, extract_strategic_initiatives, generate_opportunity_map]
PRESENTATION_TOOLS = [format_tcs_value_props, save_as_powerpoint]
REVIEWER_TOOLS = [verify_financial_facts, check_tcs_claims, generate_review_report]
